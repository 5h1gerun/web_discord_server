from __future__ import annotations

"""Automatic tagging with Gemini AI."""

from pathlib import Path
import base64
import io
import mimetypes
import os

import google.generativeai as genai
from google.generativeai import GenerationConfig
from PyPDF2 import PdfReader
from PIL import Image, UnidentifiedImageError
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


def generate_tags(file_path: Path, original_name: str | None = None) -> str:
    """Analyze the file and return comma separated tags using Gemini.

    Parameters
    ----------
    file_path:
        実際に保存されているファイルへのパス。
    original_name:
        元のファイル名。拡張子を保持していない場合に MIME 判定へ
        利用します。
    """
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        return ""

    # API キー設定＆モデル準備（温度や出力トークン数もここでまとめて指定）
    genai.configure(api_key=api_key)
    gen_cfg = GenerationConfig(
        temperature=0.2,
        max_output_tokens=256,
    )
    model = genai.GenerativeModel(
        "gemini-2.0-flash",
        generation_config=gen_cfg,
    )

    # ファイルサイズが 500MB 以上ならタグ付けをスキップ
    try:
        if file_path.stat().st_size >= 500_000_000:
            return ""
    except Exception:
        return ""

    # ファイル読み込み
    mime, _ = mimetypes.guess_type(original_name or str(file_path))
    try:
        data = file_path.read_bytes()
    except Exception:
        return ""

    if not mime:
        mime = "application/octet-stream"

    # 1) プレーンテキスト処理
    if mime and mime.startswith("text"):
        text = data.decode(errors="ignore")
        prompt = (
            "以下のテキストから重要と思われるキーワードを5個抽出し、"
            "カンマ区切りで出力してください:\n" + text[:16000]
        )
        resp = model.generate_content(prompt)
        return resp.text.strip()

    # 2) PDF はテキスト抽出を試み、失敗時はバイナリ解析
    if mime == "application/pdf":
        try:
            reader = PdfReader(str(file_path))
            pages = [p.extract_text() or "" for p in reader.pages]
        except Exception:
            return ""
        full_text = "\n".join(pages)
        if not full_text.strip():
            b64 = base64.b64encode(data).decode()
            prompt = (
                "与えられたPDFの内容を解析し、関連するキーワードを5個"
                " 日本語で抽出してカンマ区切りで返してください。"
            )
            resp = model.generate_content([
                {"mime_type": "application/pdf", "data": b64},
                {"text": prompt},
            ])
            return resp.text.strip()
        prompt = (
            "以下のPDF本文から重要と思われるキーワードを5個抽出し、"
            "カンマ区切りで出力してください:\n" + full_text[:16000]
        )
        resp = model.generate_content(prompt)
        return resp.text.strip()

    # 2-2) Office 系ファイルはテキストへ変換して解析
    if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        try:
            doc = Document(str(file_path))
            full_text = "\n".join(p.text for p in doc.paragraphs)
        except Exception:
            return ""
        prompt = (
            "以下の文書から重要と思われるキーワードを5個抽出し、"
            "カンマ区切りで出力してください:\n" + full_text[:16000]
        )
        resp = model.generate_content(prompt)
        return resp.text.strip()

    if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        try:
            pres = Presentation(str(file_path))
            texts: list[str] = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
        except Exception:
            return ""
        full_text = "\n".join(texts)
        prompt = (
            "以下のプレゼンテーションから重要と思われるキーワードを5個抽出し、"
            "カンマ区切りで出力してください:\n" + full_text[:16000]
        )
        resp = model.generate_content(prompt)
        return resp.text.strip()

    if mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        try:
            wb = load_workbook(file_path, data_only=True)
            texts: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    texts.append(" ".join(str(c) for c in row if c is not None))
        except Exception:
            return ""
        full_text = "\n".join(texts)
        prompt = (
            "以下の表計算シートから重要と思われるキーワードを5個抽出し、"
            "カンマ区切りで出力してください:\n" + full_text[:16000]
        )
        resp = model.generate_content(prompt)
        return resp.text.strip()

    # 2-3) Gemini がサポートしていない ZIP ファイルはタグ生成をスキップ
    if mime == "application/zip":
        return ""

    # 2-4) 実行ファイルはセキュリティ上スキップ
    if mime in {"application/x-msdos-program", "application/x-msdownload"}:
        return ""

    # 3) MIME 未判定だがテキストとして解釈できる場合
    if b"\x00" not in data:
        try:
            text = data.decode()
        except UnicodeDecodeError:
            text = None
        if text and text.strip():
            prompt = (
                "以下のテキストから重要と思われるキーワードを5個抽出し、"
                "カンマ区切りで出力してください:\n" + text[:16000]
            )
            resp = model.generate_content(prompt)
            return resp.text.strip()

    # 4) 画像やその他バイナリ
    if not mime or mime == "application/octet-stream":
        # MIME が未判定の場合は画像として扱えるか試みる
        try:
            img = Image.open(file_path)
            with io.BytesIO() as buf:
                img.save(buf, format="PNG")
                data = buf.getvalue()
            mime = "image/png"
        except (UnidentifiedImageError, OSError):
            # Gemini では application/octet-stream を受け付けないため
            # 未対応ファイルはタグ生成をスキップする
            return ""

    b64 = base64.b64encode(data).decode()
    prompt = (
        f"与えられた {mime} の内容を解析し、関連するキーワードを5個"
        " 日本語で抽出してカンマ区切りで返してください。"
    )
    resp = model.generate_content([
        {
            "mime_type": mime,  # サポートされている MIME タイプのみ送信
            "data": b64,
        },
        {
            "text": prompt,  # Part として扱われるプロンプト
        },
    ])
    return resp.text.strip()
